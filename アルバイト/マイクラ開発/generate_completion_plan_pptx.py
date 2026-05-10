"""完成計画書（A4 pptx・クライアント確認用・編集可能）"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUT = Path(__file__).parent / "完成計画書_クライアント確認用_20260421.pptx"

# ── テーマカラー ──
NAVY = RGBColor(0x20, 0x3A, 0x60)
ACCENT = RGBColor(0xD2, 0x8C, 0x32)
GRAY_BG = RGBColor(0xEE, 0xF0, 0xF4)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
BORDER = RGBColor(0xBE, 0xC6, 0xD2)
TEXT = RGBColor(0x1E, 0x23, 0x2D)
MUTED = RGBColor(0x6E, 0x76, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Meiryo"

# A4縦 210×297mm
SLIDE_W = Cm(21.0)
SLIDE_H = Cm(29.7)
MARGIN = Cm(1.5)


def _set_run(run, size=10, bold=False, color=TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=10, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None,
             line_color=None, margin=Cm(0.15), leading=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = Cm(0.1)
    tf.margin_bottom = Cm(0.1)
    tf.vertical_anchor = anchor
    if fill is not None:
        fill_fmt = tb.fill
        fill_fmt.solid()
        fill_fmt.fore_color.rgb = fill
    else:
        tb.fill.background()
    if line_color is not None:
        tb.line.color.rgb = line_color
        tb.line.width = Pt(0.5)
    else:
        tb.line.fill.background()

    lines = str(text).split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        run = p.add_run()
        run.text = ln
        _set_run(run, size=size, bold=bold, color=color)
    return tb


def add_rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


def add_accent_bar(slide, x, y):
    """セクション見出しの左に置くアクセントバー"""
    add_rect(slide, x, y + Cm(0.15), Cm(0.15), Cm(0.55), ACCENT)


def add_section_title(slide, x, y, text):
    add_accent_bar(slide, x, y)
    add_text(slide, x + Cm(0.3), y, Cm(19), Cm(0.9), text,
             size=14, bold=True, color=NAVY)


def add_page_title(slide, x, y, text):
    add_text(slide, x, y, Cm(18), Cm(1.0), text, size=18, bold=True, color=NAVY)
    # 下線
    add_rect(slide, x, y + Cm(1.0), Cm(6.0), Cm(0.08), NAVY)


def add_footer(slide, page, total):
    add_text(slide, MARGIN, SLIDE_H - Cm(1.0), SLIDE_W - 2 * MARGIN, Cm(0.6),
             f"ワンライフ様 Minecraft学習コンテンツ 完成計画書    {page} / {total}",
             size=8, color=MUTED, align=PP_ALIGN.CENTER)


# ── スライド生成関数 ──
def make_pres():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank


# ── Slide 1：表紙 ──
def slide_cover(prs):
    s = blank_slide(prs)
    # ヘッダーバー
    add_rect(s, 0, 0, SLIDE_W, Cm(4.0), NAVY)
    add_text(s, MARGIN, Cm(1.0), Cm(18), Cm(0.7),
             "株式会社ワンライフ 御中",
             size=11, bold=True, color=WHITE)
    add_text(s, MARGIN, Cm(1.8), Cm(18), Cm(1.3),
             "Minecraft学習コンテンツ 完成計画書",
             size=22, bold=True, color=WHITE)
    add_text(s, MARGIN, Cm(3.0), Cm(18), Cm(0.8),
             "2026年内のプロジェクト完了に向けたスケジュール共有",
             size=12, color=WHITE)

    # 本資料について ボックス
    bx, by, bw, bh = MARGIN, Cm(5.2), SLIDE_W - 2 * MARGIN, Cm(5.2)
    add_rect(s, bx, by, bw, bh, LIGHT_BG, line=BORDER)
    add_text(s, bx + Cm(0.4), by + Cm(0.35), bw - Cm(0.8), Cm(0.7),
             "本資料について", size=12, bold=True, color=NAVY)
    msg = (
        "本資料は、放課後等デイサービス向け Minecraft学習コンテンツを\n"
        "2026年内に完成・納品するための月次スケジュール、成果物の範囲、\n"
        "およびクライアント様にご協力いただきたい事項を整理したものです。\n"
        "中核となる Minecraftワールド は 2026年8月末の完成を目標とし、\n"
        "以降はPDF教材・運用マニュアルの仕上げとテストプレイを経て\n"
        "2026年12月の最終納品を目指します。"
    )
    add_text(s, bx + Cm(0.4), by + Cm(1.1), bw - Cm(0.8), Cm(4.0), msg,
             size=10.5, color=TEXT, leading=1.35)

    # 基本情報テーブル
    rows = [
        ("プロジェクト名", "前橋市をテーマとした探求学習コンテンツ（Minecraft Education）"),
        ("委託者", "株式会社ワンライフ"),
        ("受託者", "丸岡 大也"),
        ("対象", "放課後等デイサービスに通う小学1年〜高校生"),
        ("ワールド完成目標", "2026年8月末"),
        ("プロジェクト完了目標", "2026年12月末（最終納品）"),
        ("本計画書 作成日", "2026年4月21日"),
    ]
    tx, ty = MARGIN, Cm(11.0)
    col1_w = Cm(5.5)
    col2_w = SLIDE_W - 2 * MARGIN - col1_w
    row_h = Cm(1.0)
    for i, (label, value) in enumerate(rows):
        ry = ty + i * row_h
        add_rect(s, tx, ry, col1_w, row_h, GRAY_BG, line=BORDER)
        add_rect(s, tx + col1_w, ry, col2_w, row_h, WHITE, line=BORDER)
        add_text(s, tx + Cm(0.2), ry, col1_w - Cm(0.2), row_h, label,
                 size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + col1_w + Cm(0.2), ry, col2_w - Cm(0.2), row_h, value,
                 size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    # 注記
    note_y = ty + len(rows) * row_h + Cm(0.6)
    add_text(s, MARGIN, note_y, SLIDE_W - 2 * MARGIN, Cm(2.2),
             "※ 本計画は、先日共有済みの「方針すり合わせ資料」に基づくミッション型構成を前提としています。\n"
             "※ ハブ拠点の舞台地（駅周辺にするか等）は現在検討中のため、本書では未確定として扱っています。\n"
             "※ 大幅な仕様変更が発生した場合、スケジュールは再調整いたします。",
             size=9, color=MUTED, leading=1.4)


# ── Slide 2：完成目標と全体像 ──
def slide_overview(prs, page, total):
    s = blank_slide(prs)
    add_page_title(s, MARGIN, Cm(1.2), "1. 完成目標と全体像")

    add_section_title(s, MARGIN, Cm(3.0), "最終成果物（納品物）")
    add_text(s, MARGIN, Cm(4.0), SLIDE_W - 2 * MARGIN, Cm(4.5),
             "・Minecraft 統合ワールドファイル（.mcworld） × 1\n"
             "    5カテゴリ × 3難易度 = 全15ミッションエリア + ハブ拠点を収録\n"
             "・PDF教材 × 15セット（児童向けワークシート + スタッフ向けガイド）\n"
             "・スタッフ向け運用マニュアル × 1（ワールドの開閉・バックアップ・トラブル対応）",
             size=10.5, color=TEXT, leading=1.5)

    # 月次ガント
    add_section_title(s, MARGIN, Cm(8.8), "月次スケジュール概観（2026年4月〜12月）")

    months = ["4月", "5月", "6月", "7月", "8月", "9月", "10月", "11月", "12月"]
    rows = [
        ("フェーズ0：方針確定・合意",              [1, 0, 0, 0, 0, 0, 0, 0, 0]),
        ("フェーズ1：ハブ + 第1ミッション",        [0, 1, 0, 0, 0, 0, 0, 0, 0]),
        ("フェーズ2：まち・自然 計6ミッション",    [0, 0, 1, 0, 0, 0, 0, 0, 0]),
        ("フェーズ3：もの・ことば 計6ミッション",  [0, 0, 0, 1, 0, 0, 0, 0, 0]),
        ("フェーズ4：せかい + ワールド統合",       [0, 0, 0, 0, 1, 0, 0, 0, 0]),
        ("フェーズ5：PDF教材 全量仕上げ",          [0, 0, 0, 0, 0, 1, 0, 0, 0]),
        ("フェーズ6：運用マニュアル + 内部テスト", [0, 0, 0, 0, 0, 0, 1, 0, 0]),
        ("フェーズ7：現場テストプレイ",            [0, 0, 0, 0, 0, 0, 0, 1, 0]),
        ("フェーズ8：最終修正・納品",              [0, 0, 0, 0, 0, 0, 0, 0, 1]),
    ]
    gx, gy = MARGIN, Cm(9.8)
    label_w = Cm(7.5)
    chart_w = SLIDE_W - 2 * MARGIN - label_w
    cell_w = chart_w / len(months)
    header_h = Cm(0.75)
    row_h = Cm(0.75)

    # ヘッダ
    add_rect(s, gx, gy, label_w, header_h, NAVY)
    add_text(s, gx + Cm(0.2), gy, label_w - Cm(0.2), header_h, "フェーズ",
             size=9.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for i, m in enumerate(months):
        cx = gx + label_w + i * cell_w
        add_rect(s, cx, gy, cell_w, header_h, NAVY)
        add_text(s, cx, gy, cell_w, header_h, m,
                 size=9.5, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, margin=Cm(0.05))

    # データ行
    for ri, (label, bar) in enumerate(rows):
        ry = gy + header_h + ri * row_h
        base_fill = GRAY_BG if ri % 2 == 0 else LIGHT_BG
        add_rect(s, gx, ry, label_w, row_h, base_fill, line=BORDER)
        add_text(s, gx + Cm(0.2), ry, label_w - Cm(0.2), row_h, label,
                 size=9, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        for i, v in enumerate(bar):
            cx = gx + label_w + i * cell_w
            if v:
                add_rect(s, cx, ry, cell_w, row_h, ACCENT, line=BORDER)
            else:
                add_rect(s, cx, ry, cell_w, row_h, base_fill, line=BORDER)

    # 凡例
    legend_y = gy + header_h + len(rows) * row_h + Cm(0.4)
    add_rect(s, gx, legend_y + Cm(0.12), Cm(0.6), Cm(0.4), ACCENT)
    add_text(s, gx + Cm(0.8), legend_y, Cm(16), Cm(0.7),
             "＝ 当該フェーズの主担当期間（前後フェーズとは一部オーバーラップして進行）",
             size=9, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

    # 区切り
    sep_y = legend_y + Cm(1.0)
    add_section_title(s, MARGIN, sep_y, "2つのマイルストーン")
    add_rect(s, MARGIN, sep_y + Cm(1.0), (SLIDE_W - 2 * MARGIN - Cm(0.4)) / 2, Cm(3.0),
             LIGHT_BG, line=BORDER)
    add_rect(s, MARGIN + (SLIDE_W - 2 * MARGIN - Cm(0.4)) / 2 + Cm(0.4), sep_y + Cm(1.0),
             (SLIDE_W - 2 * MARGIN - Cm(0.4)) / 2, Cm(3.0), LIGHT_BG, line=BORDER)

    half_w = (SLIDE_W - 2 * MARGIN - Cm(0.4)) / 2
    # 左：ワールド完成
    add_text(s, MARGIN + Cm(0.3), sep_y + Cm(1.1), half_w - Cm(0.6), Cm(0.7),
             "① ワールド完成（2026年8月末）",
             size=11, bold=True, color=NAVY)
    add_text(s, MARGIN + Cm(0.3), sep_y + Cm(1.9), half_w - Cm(0.6), Cm(2.0),
             "全15ミッション + ハブ拠点を収録した .mcworld を完成。\n"
             "クライアント様とのレビューを経て、ワールド単体の動作を確定。",
             size=10, color=TEXT, leading=1.4)
    # 右：プロジェクト完了
    rx = MARGIN + half_w + Cm(0.4)
    add_text(s, rx + Cm(0.3), sep_y + Cm(1.1), half_w - Cm(0.6), Cm(0.7),
             "② プロジェクト完了（2026年12月末）",
             size=11, bold=True, color=NAVY)
    add_text(s, rx + Cm(0.3), sep_y + Cm(1.9), half_w - Cm(0.6), Cm(2.0),
             "PDF教材15セット + 運用マニュアル の仕上げと、\n"
             "スタッフ・児童によるテストプレイの反映を経て最終納品。",
             size=10, color=TEXT, leading=1.4)

    add_footer(s, page, total)


# ── Slide 3 & 4：フェーズ別詳細 ──
PHASES = [
    {
        "title": "フェーズ0｜方針確定・クライアント合意",
        "period": "2026年4月21日 〜 4月30日（約1.5週）",
        "goal": "ミッション型構成・スコープについて合意し、開発着手の前提を揃える。",
        "tasks": [
            "本計画書・方針すり合わせ資料に基づくご確認と承認取得",
            "ミッション名称（15本）の最終確定",
            "MakeCodeの位置づけ（段階導入／必須／お任せ）の決定",
            "「せかいをのぞこう（海外都市）」の今回スコープ可否の決定",
            "ハブ拠点の舞台地コンセプトの決定",
        ],
        "deliv": "ご承認済み要件定義書（第6版）",
    },
    {
        "title": "フェーズ1｜ハブ拠点 + 第1ミッション（品質基準の確立）",
        "period": "2026年5月1日 〜 5月31日（約4週）",
        "goal": "完成形の型を1本作り、以降の横展開用テンプレートを固める。",
        "tasks": [
            "ハブ拠点：チュートリアル神殿／5ゲート／マイルーム／テレポート機構",
            "「まちをつくろう／やさしい：公園に花壇を設計しよう」完成",
            "探求学習9ステップを全実装し、動作確認",
            "PDF教材テンプレート策定 + 第1号ドラフト作成",
            "試作版をクライアント様に共有・フィードバック取得",
        ],
        "deliv": "試作ワールド（ハブ + 第1ミッション）、PDF教材テンプレ + サンプル1セット",
    },
    {
        "title": "フェーズ2｜「まちをつくろう」「自然を探ろう」完成",
        "period": "2026年6月1日 〜 6月30日（約4週）",
        "goal": "2カテゴリ×3難易度（計6ミッション）を完成させ、ワールドの半分を仕上げる。",
        "tasks": [
            "まちをつくろう（ふつう：駅前広場リデザイン／むずかしい：商店街再開発）",
            "自然を探ろう（易：赤城山観察／中：利根川の洪水対策／上：気候と農業）",
            "各ミッションのPDF教材ドラフト（計6セット）作成",
            "ハブ⇔各ミッションのテレポート・戻り動線 全体テスト",
        ],
        "deliv": "6ミッション完成版、PDF教材ドラフト 6セット",
    },
    {
        "title": "フェーズ3｜「ものをつくろう」「ことばで伝えよう」完成",
        "period": "2026年7月1日 〜 7月31日（約4週）",
        "goal": "残り2カテゴリを完成させ、国内向け12ミッションを揃える。",
        "tasks": [
            "ものをつくろう（易：木の橋／中：耐震住宅／上：太陽光発電建物）",
            "ことばで伝えよう（易：案内看板／中：観光マップ／上：市長への提案書）",
            "各ミッションのPDF教材ドラフト（計6セット）作成",
        ],
        "deliv": "6ミッション完成版、PDF教材ドラフト 6セット",
    },
    {
        "title": "フェーズ4｜「せかいをのぞこう」完成 + ワールド統合 ★ワールド完成",
        "period": "2026年8月1日 〜 8月31日（約4週）",
        "goal": "最後のカテゴリを仕上げ、全ミッションを統合した .mcworld を完成させる。",
        "tasks": [
            "せかいをのぞこう（易：他都市比較／中：ロンドン／上：NYC × 前橋）",
            "PDF教材ドラフト 残り3セット作成",
            "全15ミッションの統合テスト、.mcworld ファイルサイズ最終確認",
            "クライアント様にワールド完成版をご報告・中間レビュー",
        ],
        "deliv": "★ Minecraftワールド完成版（.mcworld × 1、全15ミッション収録）",
    },
    {
        "title": "フェーズ5｜PDF教材 全量仕上げ・校正",
        "period": "2026年9月1日 〜 9月30日（約4週）",
        "goal": "ドラフト状態のPDF教材を製本レベルまで仕上げる。",
        "tasks": [
            "PDF教材 全15セットの本文校正・体裁統一（児童用・スタッフ用）",
            "挿絵・スクリーンショット差し替え、用語統一",
            "誤字脱字・発達特性配慮表現のチェック",
            "クライアント様にPDF一式をご確認依頼",
        ],
        "deliv": "PDF教材 全15セット（最終版候補）",
    },
    {
        "title": "フェーズ6｜スタッフ向け運用マニュアル + 内部テスト",
        "period": "2026年10月1日 〜 10月31日（約4週）",
        "goal": "現場で「ワールドを開くだけで運用できる」状態に仕上げる。",
        "tasks": [
            "スタッフ向け運用マニュアル作成（開閉・バックアップ・トラブル対応）",
            "OneDriveバックアップ運用の動作確認",
            "受託者側での操作テスト（ハブ〜全ミッション〜戻り動線）",
            "フィードバック反映・軽微な不具合修正",
        ],
        "deliv": "運用マニュアル ver.1、バグ修正済みワールド",
    },
    {
        "title": "フェーズ7｜現場テストプレイ（スタッフ・児童）",
        "period": "2026年11月1日 〜 11月30日（約4週）",
        "goal": "実際の運用環境で試し、詰まる箇所や改善点を洗い出す。",
        "tasks": [
            "スタッフ様による操作テスト（マニュアルだけで起動できるか確認）",
            "児童によるテストプレイ（9ステップが機能しているか）",
            "詰まる箇所・ストレスポイントの抽出と記録",
            "優先度付けした修正リストの作成",
        ],
        "deliv": "テスト結果レポート、修正優先リスト",
    },
    {
        "title": "フェーズ8｜最終修正・最終納品 ★プロジェクト完了",
        "period": "2026年12月1日 〜 12月下旬（約3〜4週）",
        "goal": "テスト結果を反映し、安心して現場で使える状態で納品する。",
        "tasks": [
            "テスト結果に基づく修正（ワールド・PDF・マニュアル）",
            "最終版 .mcworld のファイルサイズ・動作確認",
            "納品パッケージの整備（ファイル構成・README）",
            "最終納品とお引き渡し、納品後のフォロー方針確定",
        ],
        "deliv": "★ 最終納品パッケージ一式（.mcworld × 1 / PDF × 15 / 運用マニュアル × 1）",
    },
]


def phase_card(slide, x, y, w, phase):
    """1フェーズ分の情報カード。戻り値＝使用した高さ"""
    # タイトル帯
    th = Cm(0.85)
    add_rect(slide, x, y, w, th, NAVY)
    add_text(slide, x + Cm(0.25), y, w - Cm(0.5), th, phase["title"],
             size=10.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # 本体
    body_y = y + th
    label_w = Cm(2.1)
    val_w = w - label_w

    rows = [
        ("期間", phase["period"], 1),
        ("目的", phase["goal"], 1),
        ("主タスク", "\n".join("・" + t for t in phase["tasks"]), len(phase["tasks"])),
        ("成果物", phase["deliv"], 1),
    ]

    cy = body_y
    for label, value, lines in rows:
        # 行高さの概算
        if label == "主タスク":
            rh = Cm(0.5) * lines + Cm(0.3)
        else:
            rh = Cm(0.75)

        # 左ラベル
        add_rect(slide, x, cy, label_w, rh, GRAY_BG, line=BORDER)
        add_text(slide, x + Cm(0.15), cy, label_w - Cm(0.2), rh, label,
                 size=9.5, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        # 右値
        val_fill = LIGHT_BG if label == "成果物" else WHITE
        val_color = NAVY if label == "成果物" else TEXT
        val_bold = True if label == "成果物" else False
        add_rect(slide, x + label_w, cy, val_w, rh, val_fill, line=BORDER)
        add_text(slide, x + label_w + Cm(0.2), cy, val_w - Cm(0.3), rh, value,
                 size=9.5, color=val_color, bold=val_bold,
                 anchor=MSO_ANCHOR.MIDDLE if label != "主タスク" else MSO_ANCHOR.TOP,
                 leading=1.3)
        cy += rh

    return (cy - y)


def slide_phases_a(prs, page, total):
    s = blank_slide(prs)
    add_page_title(s, MARGIN, Cm(1.2), "2. フェーズ別 詳細スケジュール（前半：4月〜8月）")

    y = Cm(3.0)
    w = SLIDE_W - 2 * MARGIN
    for p in PHASES[:5]:
        h = phase_card(s, MARGIN, y, w, p)
        y += h + Cm(0.25)

    add_footer(s, page, total)


def slide_phases_b(prs, page, total):
    s = blank_slide(prs)
    add_page_title(s, MARGIN, Cm(1.2), "2. フェーズ別 詳細スケジュール（後半：9月〜12月）")

    y = Cm(3.0)
    w = SLIDE_W - 2 * MARGIN
    for p in PHASES[5:]:
        h = phase_card(s, MARGIN, y, w, p)
        y += h + Cm(0.3)

    add_footer(s, page, total)


# ── Slide 5：チェックポイント ──
def slide_checkpoints(prs, page, total):
    s = blank_slide(prs)
    add_page_title(s, MARGIN, Cm(1.2), "3. クライアント様とのチェックポイント")

    add_text(s, MARGIN, Cm(2.8), SLIDE_W - 2 * MARGIN, Cm(1.5),
             "各フェーズ末に、クライアント様へ進捗報告と成果物の共有を行います。\n"
             "確認いただく内容と想定スケジュールは以下の通りです。",
             size=10.5, color=TEXT, leading=1.5)

    rows = [
        ("4月末",   "フェーズ0末", "本計画書・方針すり合わせ資料のご承認",       "メール返信 / オンライン30分"),
        ("5月末",   "フェーズ1末", "ハブ + 第1ミッション試作版レビュー",       "Minecraft体験 + 方向性確認(60分)"),
        ("6月末",   "フェーズ2末", "6ミッション + PDF教材ドラフトの中間確認",  "動画共有 + フィードバック"),
        ("7月末",   "フェーズ3末", "12ミッション完成状況のご報告",               "進捗レポート（文書）"),
        ("8月末",   "ワールド完成", "★ 全15ミッション収録ワールドの確認",        "オンライン説明(60分) + .mcworld共有"),
        ("9月末",   "フェーズ5末", "PDF教材 全15セット（最終版候補）の確認",     "PDF共有 + フィードバック"),
        ("10月末",  "フェーズ6末", "運用マニュアルの確認",                     "マニュアル共有"),
        ("11月中",  "テストプレイ", "スタッフ様・児童様によるテストご協力",      "現場立会い または リモート"),
        ("12月",    "最終納品",   "★ 最終納品パッケージのお渡し・受領確認",     "対面 または オンライン"),
    ]
    cols = ["時期", "区切り", "確認内容", "想定形式"]
    widths = [Cm(2.2), Cm(3.2), Cm(7.0), SLIDE_W - 2 * MARGIN - Cm(2.2) - Cm(3.2) - Cm(7.0)]

    gx, gy = MARGIN, Cm(4.8)
    header_h = Cm(0.8)
    row_h = Cm(0.85)

    # ヘッダ
    cx = gx
    for i, c in enumerate(cols):
        add_rect(s, cx, gy, widths[i], header_h, NAVY)
        add_text(s, cx, gy, widths[i], header_h, c,
                 size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=Cm(0.1))
        cx += widths[i]

    for ri, row in enumerate(rows):
        ry = gy + header_h + ri * row_h
        base = LIGHT_BG if ri % 2 == 0 else WHITE
        cx = gx
        for i, cell in enumerate(row):
            add_rect(s, cx, ry, widths[i], row_h, base, line=BORDER)
            align = PP_ALIGN.CENTER if i < 2 else PP_ALIGN.LEFT
            is_milestone = "★" in cell
            add_text(s, cx + Cm(0.15), ry, widths[i] - Cm(0.2), row_h, cell,
                     size=9.5, color=NAVY if is_milestone else TEXT,
                     bold=is_milestone,
                     align=align, anchor=MSO_ANCHOR.MIDDLE, margin=Cm(0.05))
            cx += widths[i]

    # 下段：クライアントにお願いしたいこと
    sy = gy + header_h + len(rows) * row_h + Cm(0.7)
    add_section_title(s, MARGIN, sy, "クライアント様にご協力いただきたいこと")
    add_text(s, MARGIN, sy + Cm(1.0), SLIDE_W - 2 * MARGIN, Cm(5.0),
             "・各チェックポイントでの1〜2営業日以内のご回答\n"
             "・フェーズ1末（5月末）試作版への率直なフィードバック\n"
             "　　→ 方向性が固まる最重要タイミングのため、特に丁寧にお願いできますと幸いです\n"
             "・11月のテストプレイ期間における現場での試遊機会のご提供\n"
             "・児童向け表現・禁則事項など、発達特性配慮に関するご助言",
             size=10.5, color=TEXT, leading=1.5)

    add_footer(s, page, total)


# ── Slide 6：リスク ──
def slide_risk(prs, page, total):
    s = blank_slide(prs)
    add_page_title(s, MARGIN, Cm(1.2), "4. リスクと対応方針")

    risks = [
        ("Arnisによる海外都市地形の生成に時間がかかる",
         "「せかいをのぞこう」をフェーズ4（8月）に配置。生成遅延時は『景観のみ』にスコープ縮退。"),
        (".mcworld のファイルサイズが肥大化し動作が重くなる",
         "各フェーズ終了時にサイズ計測。必要に応じ地形精度の段階的ダウン調整で対応。"),
        ("フェーズ1試作後に方針変更の要望が入る",
         "フェーズ1を最小プロトタイプと位置づけ、早期に方向性を固定。以降は差分調整に留める。"),
        ("PDF教材15セットの作成工数が膨らむ",
         "スタッフ用スクリプトを共通テンプレート化し、ミッション差分のみ記述する形式で効率化。"),
        ("11月のテストプレイ機会が確保できない",
         "リモート動画共有／内部代替テストでの補完を併用し、12月納期を死守。"),
        ("年末の最終納品タイミングが関係者の都合で取りづらい",
         "12月前半に最終版を完成させ、後半2週間を受け渡しバッファとして確保。"),
    ]

    gx, gy = MARGIN, Cm(3.0)
    w = SLIDE_W - 2 * MARGIN
    c1 = Cm(7.5)
    c2 = w - c1
    header_h = Cm(0.85)
    row_h = Cm(2.0)

    # ヘッダ
    add_rect(s, gx, gy, c1, header_h, NAVY)
    add_text(s, gx, gy, c1, header_h, "リスク",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, gx + c1, gy, c2, header_h, NAVY)
    add_text(s, gx + c1, gy, c2, header_h, "対応方針",
             size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    for i, (risk, resp) in enumerate(risks):
        ry = gy + header_h + i * row_h
        base = LIGHT_BG if i % 2 == 0 else WHITE
        add_rect(s, gx, ry, c1, row_h, base, line=BORDER)
        add_text(s, gx + Cm(0.2), ry, c1 - Cm(0.3), row_h, risk,
                 size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE, leading=1.3)
        add_rect(s, gx + c1, ry, c2, row_h, base, line=BORDER)
        add_text(s, gx + c1 + Cm(0.2), ry, c2 - Cm(0.3), row_h, resp,
                 size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE, leading=1.3)

    # 結び
    cy = gy + header_h + len(risks) * row_h + Cm(1.0)
    add_rect(s, gx, cy, w, Cm(2.8), LIGHT_BG, line=BORDER)
    add_text(s, gx + Cm(0.4), cy + Cm(0.3), w - Cm(0.8), Cm(2.2),
             "本計画はあくまで現時点のベストプランであり、試作フェーズでの気づきや\n"
             "クライアント様からのご要望を踏まえ、柔軟に調整してまいります。\n"
             "ご不明点・ご相談事項がございましたらお気軽にご連絡ください。",
             size=10.5, color=TEXT, leading=1.5, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def main():
    prs = make_pres()
    total = 6
    slide_cover(prs)
    slide_overview(prs, 2, total)
    slide_phases_a(prs, 3, total)
    slide_phases_b(prs, 4, total)
    slide_checkpoints(prs, 5, total)
    slide_risk(prs, 6, total)
    prs.save(str(OUT))
    print(f"wrote: {OUT}")


if __name__ == "__main__":
    main()
