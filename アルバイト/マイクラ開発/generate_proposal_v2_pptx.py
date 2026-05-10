"""提案書 第2版 pptx（17ページ・A4縦・クライアント送付用）"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

OUT = Path(__file__).parent / "提案書_ワンライフ様Minecraft学習プログラム_第2版_20260421.pptx"

NAVY = RGBColor(0x20, 0x3A, 0x60)
ACCENT = RGBColor(0xD2, 0x8C, 0x32)
ACCENT2 = RGBColor(0xA0, 0x40, 0x40)
GRAY_BG = RGBColor(0xEE, 0xF0, 0xF4)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
BORDER = RGBColor(0xBE, 0xC6, 0xD2)
TEXT = RGBColor(0x1E, 0x23, 0x2D)
MUTED = RGBColor(0x6E, 0x76, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
STAR = RGBColor(0xC9, 0x8A, 0x1B)
HIGHLIGHT = RGBColor(0xFF, 0xF5, 0xDE)

FONT = "Meiryo"
SLIDE_W = Cm(21.0)
SLIDE_H = Cm(29.7)
MARGIN = Cm(1.5)
TOTAL = 17


# ========== Helpers ==========
def _set_run(run, size=10, bold=False, color=TEXT, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=10, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None,
             line_color=None, margin=Cm(0.15), leading=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = Cm(0.08)
    tf.margin_bottom = Cm(0.08)
    tf.vertical_anchor = anchor
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if line_color is not None:
        tb.line.color.rgb = line_color
        tb.line.width = Pt(0.5)
    else:
        tb.line.fill.background()
    lines = str(text).split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        run = p.add_run()
        run.text = ln
        _set_run(run, size=size, bold=bold, color=color)
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


def add_bar(slide, x, y):
    add_rect(slide, x, y + Cm(0.15), Cm(0.15), Cm(0.55), ACCENT)


def add_section_title(slide, x, y, text):
    add_bar(slide, x, y)
    add_text(slide, x + Cm(0.3), y, Cm(19), Cm(0.9), text,
             size=13, bold=True, color=NAVY)


def add_page_title(slide, x, y, text):
    add_text(slide, x, y, Cm(18), Cm(1.0), text, size=17, bold=True, color=NAVY)
    add_rect(slide, x, y + Cm(1.0), Cm(6.0), Cm(0.08), NAVY)


def add_footer(slide, page):
    if page == 1:
        return
    add_text(slide, MARGIN, SLIDE_H - Cm(0.9), SLIDE_W - 2 * MARGIN, Cm(0.5),
             f"ワンライフ様 Minecraft学習プログラム 開発提案書 第2版    {page} / {TOTAL}",
             size=8, color=MUTED, align=PP_ALIGN.CENTER)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def table(slide, gx, gy, headers, rows, col_widths,
          header_h=Cm(0.8), row_h=Cm(0.9), header_size=9.5, cell_size=9,
          highlight_star=True):
    """汎用テーブル描画。行内★で色変え"""
    # ヘッダ
    cx = gx
    for i, h in enumerate(headers):
        add_rect(slide, cx, gy, col_widths[i], header_h, NAVY)
        add_text(slide, cx, gy, col_widths[i], header_h, h,
                 size=header_size, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=Cm(0.08))
        cx += col_widths[i]
    # 行
    for ri, row in enumerate(rows):
        ry = gy + header_h + ri * row_h
        is_star = highlight_star and any("★" in str(c) for c in row)
        base = HIGHLIGHT if is_star else (LIGHT_BG if ri % 2 == 0 else WHITE)
        cx = gx
        for i, cell in enumerate(row):
            add_rect(slide, cx, ry, col_widths[i], row_h, base, line=BORDER)
            has_star = "★" in str(cell)
            add_text(slide, cx + Cm(0.1), ry, col_widths[i] - Cm(0.15), row_h,
                     cell, size=cell_size,
                     bold=has_star,
                     color=STAR if has_star else TEXT,
                     anchor=MSO_ANCHOR.MIDDLE, margin=Cm(0.05))
            cx += col_widths[i]
    return gy + header_h + len(rows) * row_h


def make_pres():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ========== P1: 表紙 ==========
def p1(prs):
    s = blank(prs)
    add_rect(s, 0, 0, SLIDE_W, Cm(5.5), NAVY)
    add_text(s, MARGIN, Cm(1.0), Cm(18), Cm(0.8),
             "株式会社ワンライフ 御中", size=11, bold=True, color=WHITE)
    add_text(s, MARGIN, Cm(2.0), Cm(18), Cm(1.4),
             "Minecraft学習プログラム", size=24, bold=True, color=WHITE)
    add_text(s, MARGIN, Cm(3.3), Cm(18), Cm(1.0),
             "開発提案書（第2版）", size=18, bold=True, color=WHITE)
    add_text(s, MARGIN, Cm(4.4), Cm(18), Cm(0.8),
             "〜 子どもの自律的な学びへ点火するプログラム 〜",
             size=11, color=WHITE)

    bx, by, bw, bh = MARGIN, Cm(7.0), SLIDE_W - 2 * MARGIN, Cm(6.3)
    add_rect(s, bx, by, bw, bh, LIGHT_BG, line=BORDER)
    add_text(s, bx + Cm(0.4), by + Cm(0.3), bw - Cm(0.8), Cm(0.8),
             "本提案書について", size=12, bold=True, color=NAVY)
    msg = ("本提案書は、放課後等デイサービス向け Minecraft学習プログラムの開発方針と\n"
           "スケジュールをまとめたものです。初版からの大きな変更点として、\n"
           "「児童の自律的な学びを点火するゲートウェイ設計思想」を正面に据え、\n"
           "Minecraftワールドと連動Webサイト・PDF教材を一体運用する設計としています。\n"
           "ご確認の上、方針のご判断（後述2点）と、進行についてのご合意をお願い申し上げます。")
    add_text(s, bx + Cm(0.4), by + Cm(1.1), bw - Cm(0.8), Cm(5.0), msg,
             size=10.5, leading=1.4)

    rows = [
        ("プロジェクト名", "Minecraft学習プログラム（ワンライフ様専用）"),
        ("委託者", "株式会社ワンライフ"),
        ("受託者", "丸岡 大也"),
        ("対象", "放課後等デイサービス利用児童（小1〜高校生）"),
        ("ワールド完成目標", "2026年8月末"),
        ("プロジェクト完了目標", "2026年12月末（最終納品）"),
        ("本提案書 作成日", "2026年4月21日"),
        ("ご回答期限の目安", "2026年4月30日（木）"),
    ]
    tx, ty = MARGIN, Cm(14.3)
    c1w = Cm(5.5)
    c2w = SLIDE_W - 2 * MARGIN - c1w
    rh = Cm(1.0)
    for i, (lbl, val) in enumerate(rows):
        ry = ty + i * rh
        add_rect(s, tx, ry, c1w, rh, GRAY_BG, line=BORDER)
        add_rect(s, tx + c1w, ry, c2w, rh, WHITE, line=BORDER)
        add_text(s, tx + Cm(0.2), ry, c1w - Cm(0.2), rh, lbl,
                 size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx + c1w + Cm(0.2), ry, c2w - Cm(0.2), rh, val,
                 size=10, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, MARGIN, SLIDE_H - Cm(2.3), SLIDE_W - 2 * MARGIN, Cm(1.5),
             "※ 本提案書は「方針すり合わせ資料」および「完成計画書」の内容を統合・改訂した第2版です。\n"
             "※ 本資料中の「推奨」は受託者側の判断材料提供であり、クライアント様のご判断を優先します。",
             size=9, color=MUTED, leading=1.4)


# ========== P2: 真の目的 ==========
def p2(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "1. 本プロジェクトの真の目的")

    add_section_title(s, MARGIN, Cm(3.0), "私たちが作るもの＝「火付け役」（ゲートウェイ）")
    msg = ("本プログラムは、児童に完結した学習を詰め込むコンテンツではなく、\n"
           "『学ぶって楽しい』『マイクラって楽しい』と感じてもらう火付け役です。\n"
           "火がついた後の本当の学びは、児童自身が選ぶサバイバル・マルチプレイ・\n"
           "創作・YouTube再現の中で起こります。本プログラムはそこへ向かうための\n"
           "素地と自信を育てる役割を担います。")
    add_text(s, MARGIN, Cm(4.0), SLIDE_W - 2 * MARGIN, Cm(5.0), msg,
             size=11, leading=1.5, color=TEXT)

    add_section_title(s, MARGIN, Cm(9.5), "この思想を選ぶ3段の論理")
    steps = [
        ("①", "制作するコンテンツには量的な限界がある",
         "15ミッション全てを完璧に踏破させることは現実的ではないし、そこに価値を置きすぎると失敗体験に繋がる。"),
        ("②", "児童は結局、自発的に自由な遊びに進む",
         "サバイバル／マルチプレイ／YouTube動画の再現など、大人が用意した枠を超えた場所で試行錯誤と協力が起こる。"),
        ("③", "だから火付け役に徹して、自律的な学びに接続する",
         "本プログラムは「入口」として設計し、その先の自由な学びへ繋がる素地・自信・語彙を育てる。"),
    ]
    y = Cm(10.6)
    for num, title, body in steps:
        add_rect(s, MARGIN, y, Cm(1.5), Cm(2.8), ACCENT)
        add_text(s, MARGIN, y, Cm(1.5), Cm(2.8), num,
                 size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, MARGIN + Cm(1.5), y, SLIDE_W - 2 * MARGIN - Cm(1.5), Cm(2.8),
                 LIGHT_BG, line=BORDER)
        add_text(s, MARGIN + Cm(1.7), y + Cm(0.2), SLIDE_W - 2 * MARGIN - Cm(1.9),
                 Cm(0.8), title, size=12, bold=True, color=NAVY)
        add_text(s, MARGIN + Cm(1.7), y + Cm(1.0), SLIDE_W - 2 * MARGIN - Cm(1.9),
                 Cm(1.6), body, size=10, leading=1.4)
        y += Cm(3.1)

    add_section_title(s, MARGIN, y + Cm(0.2), "認知機能への波及効果（ONEHOUDAY様資料と整合）")
    add_text(s, MARGIN, y + Cm(1.2), SLIDE_W - 2 * MARGIN, Cm(1.8),
             "探求サイクル『考え・試し・結果を見て・直し・伝える』を全ミッションに組み込むことで、\n"
             "高次認知機能（組織化／計画／実行機能／認知の柔軟性）の発達を促します。",
             size=10.5, leading=1.5)

    add_footer(s, 2)


# ========== P3: 前提・仮定 ==========
def p3(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "2. 本提案の前提・仮定")

    add_text(s, MARGIN, Cm(3.0), SLIDE_W - 2 * MARGIN, Cm(1.5),
             "クライアント様側でも現時点で確定していない項目がいくつかあるため、\n"
             "受託者側で設計の土台となる前提を仮置きしています。フェーズ1試作後に調整可能です。",
             size=10, color=MUTED, leading=1.4)

    add_section_title(s, MARGIN, Cm(5.0), "1. 対象層の内訳仮定")
    table(s, MARGIN, Cm(6.0),
          ["学年区分", "想定割合", "主に扱う難易度"],
          [
              ("小学校 低学年（小1〜2）", "30%", "易 中心"),
              ("小学校 中学年（小3〜4）", "40%", "易〜中"),
              ("小学校 高学年（小5〜6）", "20%", "中〜上"),
              ("中学・高校", "10%", "上＋自由課題"),
          ],
          [Cm(7.5), Cm(4.0), Cm(6.5)], row_h=Cm(0.75))

    add_section_title(s, MARGIN, Cm(10.2), "2. 1セッションの時間想定 / 3. 1ミッションあたり問題数")
    table(s, MARGIN, Cm(11.2),
          ["区分", "問題数", "役割"],
          [
              ("ウォームアップ", "2〜3問", "概念の導入・成功体験"),
              ("コア問題", "4〜5問", "学習題材の中核反復"),
              ("応用・発展", "2〜3問", "試行錯誤・自由課題への予告"),
              ("合計", "約10問", "1問あたり体感2〜4分 / 全15ミッションで約150問"),
          ],
          [Cm(4.5), Cm(3.5), Cm(10.0)], row_h=Cm(0.75))

    add_section_title(s, MARGIN, Cm(15.3), "4. 難易度設計の原則")
    table(s, MARGIN, Cm(16.3),
          ["難易度", "クリア率目安", "設計方針"],
          [
              ("易", "95%以上", "ほぼ失敗なし／自己肯定感最優先"),
              ("中", "70〜85%", "2〜3問は試行錯誤を経験"),
              ("上", "50〜70%", "3〜4問は応用・創造を要する"),
          ],
          [Cm(3.0), Cm(4.0), Cm(11.0)], row_h=Cm(0.75))

    add_section_title(s, MARGIN, Cm(20.0), "5. 調整方針")
    add_text(s, MARGIN, Cm(21.0), SLIDE_W - 2 * MARGIN, Cm(4.5),
             "・上記はあくまで設計の土台であり、確定事項ではありません。\n"
             "・フェーズ1試作版（2026年5月末完成予定）で実際に試作し、現場フィードバックを踏まえて調整します。\n"
             "・早期のミニマムプロトタイプでテスト → 全体量産前に数字を確定 → 実装リスクを抑える進行を取ります。\n"
             "・施設の実際の学年内訳・滞在時間等をご共有いただけた場合、即座に反映いたします。",
             size=10, leading=1.5)

    add_footer(s, 3)


# ========== P4: 対象と4つの提供物 ==========
def p4(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "3. 対象と4つの提供物")

    add_section_title(s, MARGIN, Cm(3.0), "対象児童・スタッフ")
    add_text(s, MARGIN, Cm(4.0), SLIDE_W - 2 * MARGIN, Cm(2.8),
             "・児童：放課後等デイサービスを利用する小学1年〜高校生（発達特性に配慮）\n"
             "・スタッフ：主に保育士（ITリテラシーゼロを前提とした運用設計）\n"
             "・環境：Minecraft Education Edition（契約済）／支給予定のマイクラ用PC",
             size=11, leading=1.6)

    add_section_title(s, MARGIN, Cm(7.2), "提供物（納品物）")
    items = [
        ("①", "Minecraft\n統合ワールド",
         ".mcworld × 1ファイル\nハブ + 4教科エリア\n+ 物理ラボ\n+ 卒業エリア"),
        ("②", "連動Webサイト",
         "ブラウザで動くミニゲーム集\n15本のマイクラ連動ゲーム\n反復ドリル・スコア記録"),
        ("③", "PDF教材 × 15セット",
         "児童用ワークシート\n+ スタッフ向けガイド\n（Web・マイクラ詳細を記載）"),
        ("④", "スタッフ運用\nマニュアル",
         "ワールドの開閉方法\nバックアップ手順\nトラブル対応"),
    ]
    item_y = Cm(8.3)
    item_w = (SLIDE_W - 2 * MARGIN - Cm(0.6)) / 4
    for i, (num, title, body) in enumerate(items):
        x = MARGIN + i * (item_w + Cm(0.2))
        add_rect(s, x, item_y, item_w, Cm(0.9), NAVY)
        add_text(s, x, item_y, item_w, Cm(0.9), num + " " + title,
                 size=10.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x, item_y + Cm(0.9), item_w, Cm(4.5), LIGHT_BG, line=BORDER)
        add_text(s, x + Cm(0.1), item_y + Cm(1.0), item_w - Cm(0.2), Cm(4.3),
                 body, size=9.5, leading=1.5,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    add_section_title(s, MARGIN, Cm(14.5), "4つの関係：ゲートウェイ構造")
    flow_y = Cm(15.8)
    flow_items = [
        ("Web\nサイト", "軽く触る\n反復ドリル"),
        ("PDF\n教材", "詳細・攻略\n応用"),
        ("マイクラ\nワールド", "実体験\n応用課題"),
        ("卒業\nエリア", "自由遊び\nへの橋渡し"),
    ]
    fw = Cm(3.6)
    gap = Cm(0.5)
    arrow_w = Cm(0.6)
    total_w = len(flow_items) * fw + (len(flow_items) - 1) * (gap + arrow_w)
    fx = (SLIDE_W - total_w) / 2

    for i, (label, sub) in enumerate(flow_items):
        cx = fx + i * (fw + gap + arrow_w)
        color = [ACCENT, NAVY, ACCENT, NAVY][i]
        add_rect(s, cx, flow_y, fw, Cm(2.6), color)
        add_text(s, cx, flow_y + Cm(0.2), fw, Cm(1.2), label,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx, flow_y + Cm(1.4), fw, Cm(1.0), sub,
                 size=9, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        if i < len(flow_items) - 1:
            arrow_x = cx + fw + Cm(0.1)
            add_rect(s, arrow_x, flow_y + Cm(1.1), arrow_w, Cm(0.4),
                     MUTED, shape=MSO_SHAPE.RIGHT_ARROW)

    add_text(s, MARGIN, flow_y + Cm(3.0), SLIDE_W - 2 * MARGIN, Cm(2.5),
             "Web＝入口で軽く、PDF＝詳細で深堀り、マイクラ＝実体験、卒業エリア＝自由遊びへの橋渡し。\n"
             "この4段階の連携が本プログラムの基本構造です。",
             size=10, color=MUTED, align=PP_ALIGN.CENTER, leading=1.5)

    add_footer(s, 4)


# ========== P5: 独自性＋競合優位性 ==========
def p5(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "4. 独自性の担保構造と競合優位性")

    add_section_title(s, MARGIN, Cm(3.0), "独自性の担保構造（コピー耐性のある設計）")
    add_text(s, MARGIN, Cm(4.0), SLIDE_W - 2 * MARGIN, Cm(2.5),
             "・Webサイトには「簡単なゲームルールのみ」を掲載\n"
             "・PDF教材には「詳細な攻略法・応用・スタッフ指導法」を掲載\n"
             "・この組み合わせにより、Webだけ見ても完全な運用は再現不可能\n"
             "・Minecraftワールド × PDF教材 × Webサイトの3点セットで初めて完全体",
             size=11, leading=1.6)

    add_section_title(s, MARGIN, Cm(9.5), "競合他社との差別化ポイント（4本柱）")
    pillars = [
        ("①日本の放デイ文脈", "日本の放課後等デイサービス現場の運用フローに完全最適化。\n保育士が迷わず運用できるPDF連動設計。"),
        ("②発達特性への配慮", "否定しない／失敗を怖がらせない／正のフィードバック多用。\nアドベンチャーモード固定による事故・迷子の防止。"),
        ("③物理ラボという独自体験", "算数×理科の交差領域で、トロッコを通じて位置エネルギー・運動エネルギー・\nレール材料計算を体感。Minecraftでしか成立しない独自コンテンツ。"),
        ("④Web+PDFの連動設計", "反復ドリルはWeb、応用はPDFで差別化。\nPDF所持者のみ完全運用可という独自性担保メカニズム。"),
    ]
    y = Cm(10.5)
    for title, body in pillars:
        add_rect(s, MARGIN, y, Cm(5.0), Cm(2.3), NAVY)
        add_text(s, MARGIN + Cm(0.2), y, Cm(4.8), Cm(2.3), title,
                 size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, MARGIN + Cm(5.0), y, SLIDE_W - 2 * MARGIN - Cm(5.0), Cm(2.3),
                 LIGHT_BG, line=BORDER)
        add_text(s, MARGIN + Cm(5.2), y + Cm(0.15), SLIDE_W - 2 * MARGIN - Cm(5.4),
                 Cm(2.0), body, size=10, leading=1.5,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += Cm(2.5)

    add_footer(s, 5)


# ========== P6: 2択① ==========
def p6(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "5.【2択①】ワールド構造の主軸")

    add_text(s, MARGIN, Cm(3.0), SLIDE_W - 2 * MARGIN, Cm(1.3),
             "ワールドをどう区切るか、2つの案を提示します。運用の手堅さから案Aを推奨します。",
             size=10.5, color=MUTED)

    half_w = (SLIDE_W - 2 * MARGIN - Cm(0.4)) / 2
    y = Cm(5.0)

    # 案A
    add_rect(s, MARGIN, y, half_w, Cm(1.2), STAR)
    add_text(s, MARGIN, y, half_w, Cm(1.2), "★ 案A：教科別エリア（推奨）",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, MARGIN, y + Cm(1.2), half_w, Cm(9.5), HIGHLIGHT, line=BORDER)
    add_text(s, MARGIN + Cm(0.3), y + Cm(1.4), half_w - Cm(0.6), Cm(9.1),
             "【構成】\n"
             "4教科 × 3難易度 = 12ミッション\n"
             "+ 物理ラボ（算数×理科交差）3ミッション\n"
             "+ 国語は全ミッションに素地統合\n"
             "+ 英語は発展接続で代替\n"
             "＝ 計15ミッション\n\n"
             "【長所】\n"
             "・スタッフの指導計画に直結\n"
             "・「今日は算数」と指示しやすい\n"
             "・学校学習との連想が強い\n"
             "・保護者への説明が容易\n\n"
             "【短所】\n"
             "・教科ラベルが前に出る\n"
             "・「勉強は道具」の実感がやや薄れる",
             size=10, leading=1.5)

    # 案B
    bx = MARGIN + half_w + Cm(0.4)
    add_rect(s, bx, y, half_w, Cm(1.2), NAVY)
    add_text(s, bx, y, half_w, Cm(1.2), "案B：従来の15ワールド方針",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, bx, y + Cm(1.2), half_w, Cm(9.5), LIGHT_BG, line=BORDER)
    add_text(s, bx + Cm(0.3), y + Cm(1.4), half_w - Cm(0.6), Cm(9.1),
             "【構成】\n"
             "5教科 × 3難易度 = 独立ワールド15個\n"
             "（国語／算数／理科／社会／英語\n"
             "  × やさしい／ふつう／むずかしい）\n"
             ".mcworld を15ファイル個別納品\n"
             "ブリーフィング原案そのままの構成\n\n"
             "【長所】\n"
             "・ブリーフィング要件に完全準拠\n"
             "・教科単位で完全に独立（混線なし）\n"
             "・「今日はこのワールド」と指示が明快\n"
             "・ワールド単位でバージョン管理しやすい\n\n"
             "【短所】\n"
             "・スタッフがワールドを15個開閉管理\n"
             "・ハブ拠点がなく世界観の没入感が弱い\n"
             "・物理ラボ等の交差領域を配置しにくい\n"
             "・卒業エリア（発展接続）が設けづらい",
             size=10, leading=1.5)

    add_section_title(s, MARGIN, Cm(15.8), "ブリーフィングとの関係")
    add_text(s, MARGIN, Cm(16.8), SLIDE_W - 2 * MARGIN, Cm(6.0),
             "ブリーフィングは「国語／算数／理科／社会／英語」の5教科網羅を要求されていました。\n"
             "本提案（案A）では、発達特性配慮とマイクラ世界観との整合を優先し、以下の構成に修正します：\n\n"
             "・国語を独立エリアに            → 全ミッションに素地統合（NPCセリフ・看板・クエスト文）\n"
             "・英語を独立エリアに            → 発展接続で代替（卒業エリアにNASA/Cambridge誘導）\n"
             "・プログラミングを全科目基礎に → 独立エリアに昇格（MakeCodeを体系的に学べる）\n"
             "・5教科構成                    → 4教科 + 物理ラボ + 国語素地 + 英語発展接続\n\n"
             "この修正は「運用負荷の軽減」と「発達特性児童への認知負荷の最適化」が主目的です。",
             size=10, leading=1.5)

    add_footer(s, 6)


# ========== P7: 2択② ==========
def p7(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "6.【2択②】独自価値の方向性")

    add_text(s, MARGIN, Cm(3.0), SLIDE_W - 2 * MARGIN, Cm(1.3),
             "プログラムの位置付けをどう定義するか。ゲートウェイ哲学と整合する案Bを推奨します。",
             size=10.5, color=MUTED)

    half_w = (SLIDE_W - 2 * MARGIN - Cm(0.4)) / 2
    y = Cm(5.0)

    # 案A
    add_rect(s, MARGIN, y, half_w, Cm(1.2), NAVY)
    add_text(s, MARGIN, y, half_w, Cm(1.2), "案A：放デイ特化・手堅く",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, MARGIN, y + Cm(1.2), half_w, Cm(9.5), LIGHT_BG, line=BORDER)
    add_text(s, MARGIN + Cm(0.3), y + Cm(1.4), half_w - Cm(0.6), Cm(9.1),
             "【位置付け】\n"
             "放デイ現場で完結する自己充足型パッケージ\n\n"
             "【長所】\n"
             "・納品後の完結性が高い\n"
             "・期待値がシンプル\n"
             "・スコープ管理が容易\n\n"
             "【短所】\n"
             "・ゲートウェイ哲学との親和性が低い\n"
             "・児童の「次」への接続が弱い\n"
             "・差別化の物語が弱い",
             size=10, leading=1.5)

    # 案B
    bx = MARGIN + half_w + Cm(0.4)
    add_rect(s, bx, y, half_w, Cm(1.2), STAR)
    add_text(s, bx, y, half_w, Cm(1.2), "★ 案B：発展接続型（推奨）",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, bx, y + Cm(1.2), half_w, Cm(9.5), HIGHLIGHT, line=BORDER)
    add_text(s, bx + Cm(0.3), y + Cm(1.4), half_w - Cm(0.6), Cm(9.1),
             "【位置付け】\n"
             "自律的な学びの入口として機能する、\n"
             "「ここを入口に、次の世界へ」という物語を持つプログラム\n\n"
             "【長所】\n"
             "・ゲートウェイ哲学と完全に整合\n"
             "・NASA／ケンブリッジ等への橋渡しを明示\n"
             "・サバイバル自由遊び・YouTube再現への接続\n"
             "・物語性のある差別化が可能\n\n"
             "【短所】\n"
             "・「発展先」を説明する追加コンテンツが必要\n"
             "・卒業エリアの制作工数がやや増える",
             size=10, leading=1.5)

    add_section_title(s, MARGIN, Cm(15.8), "発展接続の具体像")
    add_text(s, MARGIN, Cm(16.8), SLIDE_W - 2 * MARGIN, Cm(6.0),
             "・卒業エリアに、NASA の Minecraft Education ワールド／ケンブリッジ大学の公開ワールド等への\n"
             "  誘導看板・紹介を配置（英語表記／ポスター／QRコード）\n"
             "・PDF巻末に「さらに挑戦したい人へ」のページを設置：\n"
             "    - 海外Minecraft教育コンテンツの紹介\n"
             "    - おすすめYouTubeクリエイターの紹介\n"
             "    - 自分のサバイバルワールドを作るためのヒント\n"
             "・これにより、運用の手堅さは担保しつつ、『次の世界』を可視化できます。",
             size=10, leading=1.5)

    add_footer(s, 7)


# ========== P8: 設計原則 ==========
def p8(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "7. 設計原則（全ミッション共通）")

    add_section_title(s, MARGIN, Cm(3.0), "探求学習サイクル：考え・試し・結果を見て・直し・伝える")
    add_text(s, MARGIN, Cm(4.0), SLIDE_W - 2 * MARGIN, Cm(1.8),
             "ブリーフィング／ONEHOUDAY様資料の核心である探求サイクルを、\n"
             "全ミッションの基本フローとして採用します。",
             size=11, leading=1.6)

    add_section_title(s, MARGIN, Cm(6.5), "前面で謳う7つの素地スキル")
    skills = [
        ("マイクラ基本操作の自信", "反復による操作習熟。自由遊びで詰まらない土台。"),
        ("試行錯誤の習慣", "失敗を怖がらず挑戦する姿勢。否定しない設計で育てる。"),
        ("空間認識・方向感覚", "マイクラ世界での地図・座標・方向の感覚。"),
        ("相手に伝える／伝わる力", "NPC・看板・選択肢を通じて伝達の基礎を体験。"),
        ("リソース管理感覚", "必要数を先に計算する習慣。計画性の基礎。"),
        ("ルールを守って遊ぶ", "時間・エリア・行動の制約の中で工夫する姿勢。"),
        ("協調・役割意識", "1人プレイでもNPCとの協働を通じてチーム感覚を育成。"),
    ]
    y = Cm(7.5)
    col_w = (SLIDE_W - 2 * MARGIN - Cm(0.3)) / 2
    for i, (t, b) in enumerate(skills):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (col_w + Cm(0.3))
        cy = y + row * Cm(2.0)
        add_rect(s, x, cy, Cm(0.4), Cm(1.8), ACCENT)
        add_rect(s, x + Cm(0.4), cy, col_w - Cm(0.4), Cm(1.8), LIGHT_BG, line=BORDER)
        add_text(s, x + Cm(0.6), cy + Cm(0.1), col_w - Cm(0.8), Cm(0.7), t,
                 size=10.5, bold=True, color=NAVY)
        add_text(s, x + Cm(0.6), cy + Cm(0.8), col_w - Cm(0.8), Cm(1.0), b,
                 size=9, leading=1.4)

    add_section_title(s, MARGIN, Cm(16.0), "認知機能・発達特性への配慮原則")
    add_text(s, MARGIN, Cm(17.0), SLIDE_W - 2 * MARGIN, Cm(7.0),
             "・否定しない設計：失敗時にペナルティを与えない／『やり直し』がデフォルト\n"
             "・正のフィードバック多用：達成時のNPC称賛、スコア加点、エモート演出\n"
             "・アドベンチャーモード固定：事故・迷子・意図しない破壊の防止\n"
             "・1人1アカウント＋進捗自動セーブ：「続きから」で詰まる児童もペース維持\n"
             "・ひらがな多めのセリフ：低学年でも読める／漢字にはふりがな\n"
             "・視覚的な3Dアイコン多用：文字理解に時間がかかる児童もアクセス可能\n\n"
             "これらは高次認知機能（組織化・計画・実行機能・認知の柔軟性）の発達支援と整合します。",
             size=10, leading=1.5)

    add_footer(s, 8)


# ========== P9: 学習9ステップ ==========
def p9(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "8. 標準学習フロー：9ステップ")

    add_text(s, MARGIN, Cm(3.0), SLIDE_W - 2 * MARGIN, Cm(1.6),
             "ブリーフィング4.1に示された学習手順をそのまま採用し、\n"
             "全ミッションがこの9ステップに沿って進行する設計としています。",
             size=10.5, color=MUTED, leading=1.5)

    steps = [
        ("①", "操作手順の説明", "Minecraftの基本操作を説明", "未経験者への導入"),
        ("②", "操作手順の実施", "実際にキャラクターを動かす", "基本操作の習熟"),
        ("③", "簡単な課題の提供", "シンプルな課題を与え実施", "成功体験・達成感"),
        ("④", "やり方の確認（情報収集）", "課題達成に至る手順・情報を提供", "手順の理解"),
        ("⑤", "実施（整理・分析）", "学んだ手順で試行錯誤しながら取り組む", "思考力・分析力"),
        ("⑥", "結果のフィードバック", "想像と実際の結果の差異を実感", "計画と実行のズレ認識"),
        ("⑦", "完成と正のフィードバック", "完成し、スタッフが褒める", "自己肯定感の向上"),
        ("⑧", "振り返り", "満足度・うまくいかなかった点を話し合う", "内省・改善意欲"),
        ("⑨", "自由課題", "学んだ内容で本人が自由に課題設定", "主体性・創造性"),
    ]
    widths = [Cm(1.3), Cm(5.5), Cm(7.0), Cm(4.2)]
    gx, gy = MARGIN, Cm(5.2)
    # header
    cx = gx
    for i, h in enumerate(["Step", "名称", "内容", "目的"]):
        add_rect(s, cx, gy, widths[i], Cm(0.9), NAVY)
        add_text(s, cx, gy, widths[i], Cm(0.9), h,
                 size=10.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += widths[i]
    for ri, row in enumerate(steps):
        ry = gy + Cm(0.9) + ri * Cm(1.4)
        base = LIGHT_BG if ri % 2 == 0 else WHITE
        cx = gx
        for i, cell in enumerate(row):
            add_rect(s, cx, ry, widths[i], Cm(1.4), base, line=BORDER)
            if i == 0:
                add_text(s, cx, ry, widths[i], Cm(1.4), cell,
                         size=14, bold=True, color=ACCENT,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                sz = 10 if i == 1 else 9.5
                bold = (i == 1)
                add_text(s, cx + Cm(0.15), ry, widths[i] - Cm(0.2), Cm(1.4),
                         cell, size=sz, bold=bold,
                         color=NAVY if i == 1 else TEXT,
                         anchor=MSO_ANCHOR.MIDDLE, leading=1.3)
            cx += widths[i]

    add_text(s, MARGIN, Cm(18.5), SLIDE_W - 2 * MARGIN, Cm(3.0),
             "※ 全15ミッションのPDF教材は、この9ステップの構造に沿って目次が統一されています。\n"
             "※ 以降のページでは、このフローを当てはめた実際のミッション例を2本ご紹介します。",
             size=9.5, color=MUTED, leading=1.5)

    add_footer(s, 9)


# ========== P10: ウォークスルー① ==========
def p10(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "9. ウォークスルー① 算数×中「村人にパンを公平に分けよう」")

    add_section_title(s, MARGIN, Cm(3.0), "ストーリー")
    add_rect(s, MARGIN, Cm(4.0), SLIDE_W - 2 * MARGIN, Cm(2.0), HIGHLIGHT, line=BORDER)
    add_text(s, MARGIN + Cm(0.3), Cm(4.1), SLIDE_W - 2 * MARGIN - Cm(0.6), Cm(1.8),
             "パン屋のハナコさん：\n"
             "「村のお祭りで、20個のパンを5軒の家に公平に分けてほしいの。\n"
             " 1軒にいくつ配ればいいか、数えて持っていってくれる？」",
             size=10.5, leading=1.5, anchor=MSO_ANCHOR.MIDDLE)

    add_section_title(s, MARGIN, Cm(6.5), "9ステップ進行")
    table(s, MARGIN, Cm(7.5),
          ["Step", "児童の体験", "スタッフ行動"],
          [
              ("①操作", "チュートリアル看板：歩く／物を拾う／渡す（初回のみ）", "見守り"),
              ("②実施", "練習エリアで1個のパンを1軒に渡す", "「できたね」と肯定"),
              ("③課題", "ハナコが依頼／PDFに「公平って？」ミニコラム", "PDFを一緒に読む"),
              ("④確認", "看板「20個÷5軒＝？」／NPC「配ってもOK、計算もOK」", "どちらで解くか選ばせる"),
              ("⑤実施", "各家に運ぶ／不正配置は防止", "詰まったらヒント"),
              ("⑥FB", "全戸4個で笑顔エモート／足りないと家主「?」", "観察を促す"),
              ("⑦完成", "ハナコ「ありがとう！」／スコア加点／勲章アイテム", "工程を褒める"),
              ("⑧振り返り", "PDFの問い：「20÷5＝4ってどういうこと？」", "1〜2問だけ軽く"),
              ("⑨自由", "「23個を5軒に分けると？」余り発見への誘導", "児童のペースに任せる"),
          ],
          [Cm(1.5), Cm(11.0), Cm(5.5)], row_h=Cm(1.2), cell_size=9)

    add_section_title(s, MARGIN, Cm(20.0), "獲得する学び・Web連動")
    add_text(s, MARGIN, Cm(21.0), SLIDE_W - 2 * MARGIN, Cm(4.5),
             "【算数】20÷5＝4 の実感／余りのある割り算（23÷5＝4余り3）\n"
             "【素地スキル】組織化（計画的配分）・試行錯誤・資源管理\n"
             "【国語自然統合】ハナコのセリフ読解／「公平」の語彙\n"
             "【Web連動】「パン分配パズル」：反復ドリル。PDFには余り扱い・2段階分配の攻略 → Webだけでは完全攻略不可",
             size=10, leading=1.5)

    add_footer(s, 10)


# ========== P11: ウォークスルー② 物理ラボ ==========
def p11(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "10. ウォークスルー② 物理ラボ×中「レールは何本必要？」")

    add_section_title(s, MARGIN, Cm(3.0), "ストーリー ／ プロジェクトの目玉ミッション")
    add_rect(s, MARGIN, Cm(4.0), SLIDE_W - 2 * MARGIN, Cm(2.3), HIGHLIGHT, line=BORDER)
    add_text(s, MARGIN + Cm(0.3), Cm(4.1), SLIDE_W - 2 * MARGIN - Cm(0.6), Cm(2.1),
             "駅長NPC・タツジさん：\n"
             "「村の東駅から森の奥駅まで線路を引きたい。ちょうど20ブロック離れてる。\n"
             " レールを過不足なく敷いてくれないか？ 鉄1個からレールは6本作れるよ。」",
             size=10.5, leading=1.5, anchor=MSO_ANCHOR.MIDDLE)

    add_section_title(s, MARGIN, Cm(6.8), "10問構成例（1ミッション = 10問）")
    table(s, MARGIN, Cm(7.8),
          ["#", "問題", "種別"],
          [
              ("1", "距離10ブロック、レール何本？（答:10本）", "ウォーム"),
              ("2", "距離15ブロック、レール何本？（答:15本）", "ウォーム"),
              ("3", "距離20ブロック、レール何本？（答:20本）", "ウォーム"),
              ("4", "鉄3個、レール何本作れる？（答:18本）", "コア"),
              ("5", "距離18ブロック、鉄は何個必要？（答:3個）", "コア"),
              ("6", "鉄5個、距離何ブロックまで敷ける？（答:30）", "コア"),
              ("7", "距離20、鉄は？（20÷6＝3余り2 → 4個必要）", "コア・切上"),
              ("8", "距離25、鉄は？（25÷6＝4余り1 → 5個必要）", "コア・切上"),
              ("9", "鉄4個で24本、距離30必要。足りる？（足りない）", "応用"),
              ("10", "自由課題：鉄8個で好きなコース設計。何ブロック？", "発展"),
          ],
          [Cm(1.0), Cm(12.5), Cm(4.5)], row_h=Cm(0.85), cell_size=9)

    add_section_title(s, MARGIN, Cm(18.0), "学び・他教科への接続")
    add_text(s, MARGIN, Cm(19.0), SLIDE_W - 2 * MARGIN, Cm(6.5),
             "【算数】距離計算／割り算と切り上げ（必要数の概念）／材料計算\n"
             "【理科予告】「平らだと止まる」「坂道で加速する」→ 物理ラボ×上「エネルギー最速コース」へ接続\n"
             "【プログラミング予告】「パワードレールで自動加速」→ プログラミングエリアへ接続\n"
             "【素地スキル】組織化（計画的材料調達）／資源管理／試行錯誤／空間認識\n"
             "【国語自然統合】「過不足なく」「用意する」の語彙\n"
             "【Web連動】「レール数計算パズル」で反復ドリル。PDFには切り上げ計算・長距離分割・\n"
             "           パワードレール間隔最適化の攻略 → マイクラ内で一段上のコース設計が可能に",
             size=10, leading=1.5)

    add_footer(s, 11)


# ========== P12: 15ミッション一覧 ==========
def p12(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "11. 15ミッション一覧")

    add_text(s, MARGIN, Cm(3.0), SLIDE_W - 2 * MARGIN, Cm(1.5),
             "4教科 × 3難易度 = 12ミッション ＋ 物理ラボ 3ミッション ＝ 計15ミッション",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    rows = [
        ("算数", "動物エサ配り\n（四則演算・足し算）",
                 "パン分配\n（割り算・余り）",
                 "倉庫の体積計算\n（立体・三次元）"),
        ("理科", "元素収集\n（元素記号・基礎）",
                 "元素合成\n（化合物・H₂O等）",
                 "ポーション調合\n（反応・変化）"),
        ("社会", "地図の宝探し\n（地図記号・方角）",
                 "立地シミュレーション\n（地形・立地条件）",
                 "特産品交易ルート\n（特産品・最適化）"),
        ("プロ\nグラミング", "エージェント移動\n（逐次処理）",
                 "ループ塀\n（繰り返し・効率化）",
                 "条件分岐装置\n（条件・変数・ロジック）"),
        ("★物理ラボ\n（算数×理科）", "トロッコ直線コース\n（乗り物の基本）",
                 "レール数計算\n（距離・材料計算）",
                 "エネルギー最速コース\n（位置E・運動E）"),
    ]

    gx, gy = MARGIN, Cm(5.2)
    widths = [Cm(3.2), Cm(4.8), Cm(5.0), Cm(5.0)]
    headers = ["エリア", "易（小1〜2想定）", "中（小3〜4想定）", "上（小5〜・中高）"]

    cx = gx
    for i, h in enumerate(headers):
        add_rect(s, cx, gy, widths[i], Cm(0.9), NAVY)
        add_text(s, cx, gy, widths[i], Cm(0.9), h,
                 size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += widths[i]

    for ri, row in enumerate(rows):
        ry = gy + Cm(0.9) + ri * Cm(2.2)
        is_physics = "★" in row[0]
        base = HIGHLIGHT if is_physics else (LIGHT_BG if ri % 2 == 0 else WHITE)
        cx = gx
        for i, cell in enumerate(row):
            add_rect(s, cx, ry, widths[i], Cm(2.2), base, line=BORDER)
            if i == 0:
                add_text(s, cx, ry, widths[i], Cm(2.2), cell,
                         size=10.5, bold=True,
                         color=STAR if is_physics else NAVY,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, leading=1.3)
            else:
                add_text(s, cx + Cm(0.15), ry, widths[i] - Cm(0.2), Cm(2.2),
                         cell, size=9.5,
                         anchor=MSO_ANCHOR.MIDDLE, leading=1.4)
            cx += widths[i]

    # 卒業エリア
    y2 = gy + Cm(0.9) + len(rows) * Cm(2.2) + Cm(0.5)
    add_rect(s, MARGIN, y2, SLIDE_W - 2 * MARGIN, Cm(2.0), NAVY)
    add_text(s, MARGIN + Cm(0.4), y2 + Cm(0.15), SLIDE_W - 2 * MARGIN - Cm(0.8),
             Cm(0.7), "卒業エリア（全15ミッションクリア後に解放）",
             size=11, bold=True, color=WHITE)
    add_text(s, MARGIN + Cm(0.4), y2 + Cm(0.85), SLIDE_W - 2 * MARGIN - Cm(0.8),
             Cm(1.2),
             "・限定的な自由建築／サバイバル風体験ゾーン\n"
             "・英語看板・NASA/ケンブリッジへの誘導掲示／YouTuber紹介",
             size=9.5, color=WHITE, leading=1.4)

    add_text(s, MARGIN, SLIDE_H - Cm(2.0), SLIDE_W - 2 * MARGIN, Cm(1.2),
             "※ 国語は全ミッションのNPCセリフ・看板・クエスト文で素地として自然統合。独立エリアなし。",
             size=9.5, color=MUTED, leading=1.4)

    add_footer(s, 12)


# ========== P13: Webサイト設計 ==========
def p13(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "12. Webサイト設計（15ゲーム）")

    add_section_title(s, MARGIN, Cm(3.0), "役割分担：Webで反復、マイクラで応用")
    add_text(s, MARGIN, Cm(4.0), SLIDE_W - 2 * MARGIN, Cm(2.5),
             "・Webサイト：マイクラで扱いづらい反復ドリル・スコア競争を担当\n"
             "・ブラウザで完結、インストール不要、スタッフ負担ゼロ\n"
             "・各ゲームは10ステージ構成（1ミッションの10問と対応）",
             size=10.5, leading=1.6)

    add_section_title(s, MARGIN, Cm(8.0), "15ゲームラインナップ")
    games = [
        ("1", "算数易", "動物エサ足し算バトル"),
        ("2", "算数中", "パン分配パズル"),
        ("3", "算数上", "倉庫パッキング計算"),
        ("4", "理科易", "元素収集カード"),
        ("5", "理科中", "元素コンバイナー"),
        ("6", "理科上", "ポーション調合シミュ"),
        ("7", "社会易", "地図読みクエスト"),
        ("8", "社会中", "立地シミュレーター"),
        ("9", "社会上", "特産品交易ルートパズル"),
        ("10", "プロ易", "エージェント移動パズル"),
        ("11", "プロ中", "ループ塀ビルダー"),
        ("12", "プロ上", "条件分岐装置設計"),
        ("13", "★物理易", "トロッコレール敷設ゲーム"),
        ("14", "★物理中", "レール数計算パズル"),
        ("15", "★物理上", "エネルギー最速コース設計"),
    ]
    # 3列×5行
    col_w = (SLIDE_W - 2 * MARGIN - Cm(0.4)) / 3
    y = Cm(9.0)
    for i, (n, cat, name) in enumerate(games):
        col = i % 3
        row = i // 3
        x = MARGIN + col * (col_w + Cm(0.2))
        cy = y + row * Cm(1.2)
        is_phy = "★" in cat
        base = HIGHLIGHT if is_phy else LIGHT_BG
        add_rect(s, x, cy, col_w, Cm(1.1), base, line=BORDER)
        add_text(s, x + Cm(0.2), cy, Cm(0.8), Cm(1.1), n,
                 size=12, bold=True, color=ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Cm(1.0), cy, Cm(1.8), Cm(1.1), cat,
                 size=9, bold=True,
                 color=STAR if is_phy else NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Cm(2.8), cy, col_w - Cm(2.9), Cm(1.1), name,
                 size=9, anchor=MSO_ANCHOR.MIDDLE, leading=1.3)

    add_section_title(s, MARGIN, Cm(16.0), "独自性担保メカニズム（差別化の核）")
    add_text(s, MARGIN, Cm(17.0), SLIDE_W - 2 * MARGIN, Cm(5.0),
             "・PDFに専用アクセスコード印字／初回はコード入力でWebサイト解錠\n"
             "・各ゲームの詳細ルール・応用テクニック・得点倍化法・隠し要素は全てPDFに掲載\n"
             "・マイクラ内スコアとWebスコアをゆるやかに連携\n"
             "   （PDFに「Webで100点取るとマイクラで〇〇が使える」等のヒント）\n"
             "・Webだけ見ても再現不可、PDFを持つ事業所のみ完全運用可能",
             size=10, leading=1.5)

    add_section_title(s, MARGIN, Cm(22.5), "実装順序")
    add_text(s, MARGIN, Cm(23.3), SLIDE_W - 2 * MARGIN, Cm(2.0),
             "・5〜6月：易5本（MVP・先行実装）／7〜9月：中5本／10〜11月：上5本／12月：全体仕上げ",
             size=10, leading=1.5)

    add_footer(s, 13)


# ========== P14: 1セッション運用シナリオ ==========
def p14(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "13. 1セッション運用シナリオ")

    add_text(s, MARGIN, Cm(3.0), SLIDE_W - 2 * MARGIN, Cm(1.5),
             "「スタッフ様がPDFだけ見て運用できる」状態を実現する、60分1セッションの典型的な流れです。",
             size=10.5, color=MUTED)

    scenarios = [
        ("0〜5分", "導入",
         "児童：PCの前に座る、ログイン\nスタッフ：ワールドを開く（既に起動済でもOK）、今日の目標ミッションを確認"),
        ("5〜10分", "ハブ→教科エリア",
         "児童：ハブの扉からミッションエリアへテレポート\nスタッフ：PDFの「今日のミッション」ページを開く"),
        ("10〜45分", "ミッション体験（9ステップ）",
         "児童：NPCから課題を受け、①〜⑦を進める（約35分）\nスタッフ：PDFに書かれた声かけに沿って見守り／⑥FB・⑦称賛のタイミングで関わる"),
        ("45〜55分", "振り返り（⑧⑨）",
         "児童：ワークシート（PDF児童用）を書き込む／スコア確認\nスタッフ：PDFの「問いかけ例」から1〜2問を選んで問いかける"),
        ("55〜60分", "終了",
         "児童：ログアウト／進捗はアカウントに自動保存\nスタッフ：進捗を児童ノートに1行記録（任意）"),
    ]
    widths = [Cm(2.2), Cm(3.0), Cm(12.8)]
    gx, gy = MARGIN, Cm(5.0)
    # header
    cx = gx
    for i, h in enumerate(["時間", "場面", "児童 / スタッフ 動線"]):
        add_rect(s, cx, gy, widths[i], Cm(0.9), NAVY)
        add_text(s, cx, gy, widths[i], Cm(0.9), h,
                 size=10.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += widths[i]
    for ri, row in enumerate(scenarios):
        ry = gy + Cm(0.9) + ri * Cm(2.4)
        base = LIGHT_BG if ri % 2 == 0 else WHITE
        cx = gx
        for i, cell in enumerate(row):
            add_rect(s, cx, ry, widths[i], Cm(2.4), base, line=BORDER)
            sz = 10 if i < 2 else 9.5
            bold = (i == 1)
            color = NAVY if i == 1 else TEXT
            align = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
            add_text(s, cx + Cm(0.15), ry, widths[i] - Cm(0.2), Cm(2.4),
                     cell, size=sz, bold=bold, color=color,
                     anchor=MSO_ANCHOR.MIDDLE, align=align, leading=1.4)
            cx += widths[i]

    add_section_title(s, MARGIN, Cm(19.5), "スタッフ負担ゼロの根拠")
    add_text(s, MARGIN, Cm(20.5), SLIDE_W - 2 * MARGIN, Cm(4.0),
             "・ワールドを開く以外の操作は発生しない（ログインはクライアント側PCで事前設定）\n"
             "・スタッフが判断すべきタイミングはすべてPDFに書かれている\n"
             "・IT知識ゼロ、ゲーム未経験のスタッフでも、PDFの『問いかけ例』を読むだけで役割を果たせる\n"
             "・児童が詰まっても、アドベンチャーモードのため事故は起きない",
             size=10, leading=1.5)

    add_footer(s, 14)


# ========== P15: スケジュール ==========
def p15(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "14. スケジュール（2マイルストーン）")

    # マイルストーン
    ms_y = Cm(3.0)
    half = (SLIDE_W - 2 * MARGIN - Cm(0.4)) / 2
    for i, (label, val, color) in enumerate([
        ("★ ワールド完成", "2026年8月末", ACCENT),
        ("★ プロジェクト完了", "2026年12月末", NAVY),
    ]):
        x = MARGIN + i * (half + Cm(0.4))
        add_rect(s, x, ms_y, half, Cm(2.0), color)
        add_text(s, x, ms_y + Cm(0.1), half, Cm(0.8), label,
                 size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, ms_y + Cm(1.0), half, Cm(0.8), val,
                 size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ガント
    add_section_title(s, MARGIN, Cm(6.0), "月次スケジュール（2026年4月〜12月）")
    months = ["4月", "5月", "6月", "7月", "8月", "9月", "10月", "11月", "12月"]
    phases = [
        ("フェーズ0：方針確定・合意",              [1, 0, 0, 0, 0, 0, 0, 0, 0]),
        ("フェーズ1：ハブ + 第1ミッション試作",   [0, 1, 0, 0, 0, 0, 0, 0, 0]),
        ("フェーズ2：算数・理科 6ミッション",      [0, 0, 1, 0, 0, 0, 0, 0, 0]),
        ("フェーズ3：社会・プロ 6ミッション",      [0, 0, 0, 1, 0, 0, 0, 0, 0]),
        ("フェーズ4：物理ラボ3 + ワールド統合",   [0, 0, 0, 0, 1, 0, 0, 0, 0]),
        ("フェーズ5：PDF教材 全15セット仕上げ",   [0, 0, 0, 0, 0, 1, 0, 0, 0]),
        ("フェーズ6：運用マニュアル + 内部テスト",[0, 0, 0, 0, 0, 0, 1, 0, 0]),
        ("フェーズ7：現場テストプレイ",            [0, 0, 0, 0, 0, 0, 0, 1, 0]),
        ("フェーズ8：最終修正・納品",              [0, 0, 0, 0, 0, 0, 0, 0, 1]),
        ("Webサイト開発（並行進行）",              [0, 1, 1, 1, 1, 1, 1, 0, 1]),
    ]
    gx, gy = MARGIN, Cm(7.0)
    label_w = Cm(7.0)
    chart_w = SLIDE_W - 2 * MARGIN - label_w
    cell_w = chart_w / len(months)
    header_h = Cm(0.7)
    row_h = Cm(0.72)

    add_rect(s, gx, gy, label_w, header_h, NAVY)
    add_text(s, gx + Cm(0.2), gy, label_w - Cm(0.2), header_h, "フェーズ",
             size=9.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    for i, m in enumerate(months):
        cx = gx + label_w + i * cell_w
        add_rect(s, cx, gy, cell_w, header_h, NAVY)
        add_text(s, cx, gy, cell_w, header_h, m,
                 size=9.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=Cm(0.03))

    for ri, (label, bar) in enumerate(phases):
        ry = gy + header_h + ri * row_h
        is_web = "Webサイト" in label
        base = HIGHLIGHT if is_web else (GRAY_BG if ri % 2 == 0 else LIGHT_BG)
        fill_color = STAR if is_web else ACCENT
        add_rect(s, gx, ry, label_w, row_h, base, line=BORDER)
        add_text(s, gx + Cm(0.2), ry, label_w - Cm(0.2), row_h, label,
                 size=8.5,
                 bold=is_web,
                 color=STAR if is_web else TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)
        for i, v in enumerate(bar):
            cx = gx + label_w + i * cell_w
            if v:
                add_rect(s, cx, ry, cell_w, row_h, fill_color, line=BORDER)
            else:
                add_rect(s, cx, ry, cell_w, row_h, base, line=BORDER)

    add_section_title(s, MARGIN, Cm(16.0), "主要マイルストーン説明")
    add_text(s, MARGIN, Cm(17.0), SLIDE_W - 2 * MARGIN, Cm(7.0),
             "・5月末：第1ミッション試作版をクライアント様にご確認いただく\n"
             "  → この時点でのフィードバックを基に、残り14ミッションの設計基準を確定します\n"
             "  （最重要のフィードバックタイミング）\n\n"
             "・8月末：★ ワールド完成。全15ミッションを収録した .mcworld を納品\n"
             "・11月：現場でのテストプレイを実施し、調整事項を洗い出し\n"
             "・12月末：★ 最終納品（.mcworld + PDF15セット + 運用マニュアル + Webサイト）",
             size=10, leading=1.5)

    add_footer(s, 15)


# ========== P16: クライアントへのお願い ==========
def p16(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "15. クライアント様へのお願い")

    asks = [
        ("4月30日(目安)", "本提案書へのご判断",
         "・2択① 構造の主軸（推奨A：教科別＋物理ラボ＋国語素地）\n"
         "・2択② 独自価値の方向（推奨B：発展接続型）\n"
         "・前提・仮定ページの内容（問題数10問、難易度設計等）\n"
         "上記3点について、『そのまま承認』または『ここを調整希望』をご連絡ください。"),
        ("5月末", "試作版のフィードバック",
         "ハブ拠点 + 第1ミッションの試作版を共有します。\n"
         "「方向性として違和感がないか」「児童に合うか」をご確認ください。\n"
         "このタイミングのフィードバックが残り14ミッションの品質を決める最重要ポイントです。"),
        ("各フェーズ末", "進捗報告への目視確認",
         "月1回程度、進捗レポート・動画・試作ファイルを共有します。\n"
         "1〜2営業日以内の『受領・問題なし』のお返事をお願いできますと進行が滞りません。"),
        ("11月", "現場テストプレイへのご協力",
         "実際の児童・スタッフによるテストプレイを実施させてください。\n"
         "リモート立会い、または録画共有でも構いません。"),
        ("随時", "発達特性配慮に関するご助言",
         "児童向けの表現・禁則事項・避けるべき演出など、現場知見を\n"
         "いつでもお聞かせください（Chatwork等でお願いします）。"),
    ]

    y = Cm(3.0)
    for when, title, body in asks:
        add_rect(s, MARGIN, y, Cm(3.2), Cm(3.5), ACCENT)
        add_text(s, MARGIN, y, Cm(3.2), Cm(3.5), when,
                 size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, MARGIN + Cm(3.2), y, SLIDE_W - 2 * MARGIN - Cm(3.2), Cm(3.5),
                 LIGHT_BG, line=BORDER)
        add_text(s, MARGIN + Cm(3.4), y + Cm(0.2), SLIDE_W - 2 * MARGIN - Cm(3.6),
                 Cm(0.7), title, size=11, bold=True, color=NAVY)
        add_text(s, MARGIN + Cm(3.4), y + Cm(0.9), SLIDE_W - 2 * MARGIN - Cm(3.6),
                 Cm(2.5), body, size=9.5, leading=1.5)
        y += Cm(3.8)

    add_footer(s, 16)


# ========== P17: リスク+結び ==========
def p17(prs):
    s = blank(prs)
    add_page_title(s, MARGIN, Cm(1.2), "16. リスクと対応方針／結び")

    add_section_title(s, MARGIN, Cm(3.0), "主要リスクと対応方針")
    risks = [
        ("フェーズ1試作後に方針変更の要望が入る",
         "フェーズ1を最小プロトタイプと位置付け、5月末に方向性を固定。以降は差分調整に留める。"),
        (".mcworld のファイルサイズ肥大化",
         "各フェーズ終了時にサイズ計測。必要に応じ地形精度を段階的にダウン調整。"),
        ("Webサイトの並行開発で工数が逼迫",
         "Web易5本をMVPとして5〜6月に集中実装。中・上は各フェーズに分散配置。"),
        ("PDF教材15セットの制作工数増",
         "スタッフ向けスクリプトをテンプレート化。ミッション差分のみ記述する効率運用。"),
        ("11月のテストプレイ機会が確保できない",
         "リモート録画共有／内部代替テストで補完し、12月納期を死守。"),
        ("年末の納品タイミング調整",
         "12月前半に最終版を完成。後半2週間を受け渡しバッファとして確保。"),
    ]
    widths = [Cm(8.0), Cm(10.0)]
    gx, gy = MARGIN, Cm(4.0)
    add_rect(s, gx, gy, widths[0], Cm(0.8), ACCENT2)
    add_text(s, gx, gy, widths[0], Cm(0.8), "リスク",
             size=10.5, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, gx + widths[0], gy, widths[1], Cm(0.8), NAVY)
    add_text(s, gx + widths[0], gy, widths[1], Cm(0.8), "対応方針",
             size=10.5, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for ri, (r, resp) in enumerate(risks):
        ry = gy + Cm(0.8) + ri * Cm(2.0)
        base = LIGHT_BG if ri % 2 == 0 else WHITE
        add_rect(s, gx, ry, widths[0], Cm(2.0), base, line=BORDER)
        add_text(s, gx + Cm(0.2), ry, widths[0] - Cm(0.3), Cm(2.0), r,
                 size=10, anchor=MSO_ANCHOR.MIDDLE, leading=1.4)
        add_rect(s, gx + widths[0], ry, widths[1], Cm(2.0), base, line=BORDER)
        add_text(s, gx + widths[0] + Cm(0.2), ry, widths[1] - Cm(0.3), Cm(2.0),
                 resp, size=10, anchor=MSO_ANCHOR.MIDDLE, leading=1.4)

    add_section_title(s, MARGIN, Cm(17.5), "結びのメッセージ")
    add_rect(s, MARGIN, Cm(18.5), SLIDE_W - 2 * MARGIN, Cm(6.0),
             HIGHLIGHT, line=BORDER)
    add_text(s, MARGIN + Cm(0.5), Cm(18.8), SLIDE_W - 2 * MARGIN - Cm(1.0), Cm(5.4),
             "本プログラムは、単なる15ミッションのセットではなく、\n"
             "子どもが『学ぶって楽しい』『マイクラって楽しい』と感じ、\n"
             "その火を自律的な学びへと育てていくためのゲートウェイです。\n\n"
             "クライアント様と手を取り合いながら、現場で確かに機能する\n"
             "プログラムを2026年内に完成させてまいります。\n\n"
             "どうぞよろしくお願いいたします。",
             size=11.5, leading=1.6, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER)

    add_footer(s, 17)


# ========== main ==========
def main():
    prs = make_pres()
    for fn in [p1, p2, p3, p4, p5, p6, p7, p8, p9, p10, p11, p12, p13, p14, p15, p16, p17]:
        fn(prs)
    prs.save(str(OUT))
    print(f"wrote: {OUT}")


if __name__ == "__main__":
    main()
